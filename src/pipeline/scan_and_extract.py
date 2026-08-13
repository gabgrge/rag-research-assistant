from __future__ import annotations

import argparse
import hashlib
import json
import os
from datetime import datetime, timezone
from pathlib import Path
from typing import Dict, List, Optional, Tuple

import pymupdf
from docx import Document
from dotenv import load_dotenv
from pptx import Presentation
from tqdm import tqdm

from src.utils.common_utils import utc_now_iso
from src.utils.fs_utils import path_is_under
from src.utils.logging_utils import build_script_logger
from src.utils.paths import CONVERTED_DIR, EXTRACTED_DIR, LOGS_DIR, REGISTRY_DIR
from src.core.registry import (
    EXTRACTION_EXTRACTED,
    EXTRACTION_FAILED,
    EXTRACTION_MISSING,
    EXTRACTION_NEW,
    Row,
    clear_chunk_tracking,
    load_registry_rows,
    set_chunk_pending,
    write_registry_rows,
)

load_dotenv()

_raw_dir_value = os.getenv("RAW_DIR", "").strip()
if not _raw_dir_value:
    raise FileNotFoundError("RAW_DIR not set. Define it in .env.")

RAW_DIR = Path(_raw_dir_value).expanduser().resolve()
if not RAW_DIR.exists():
    raise FileNotFoundError(f"RAW_DIR not found: {RAW_DIR}")

CONVERTED_DIR = CONVERTED_DIR.resolve()
EXTRACTED_DIR = EXTRACTED_DIR.resolve()
REGISTRY_PATH = (REGISTRY_DIR / "sources.csv").resolve()
LOG_PATH = (LOGS_DIR / "scan_and_extract.log").resolve()

EXTRACTED_DIR.mkdir(parents=True, exist_ok=True)
REGISTRY_PATH.parent.mkdir(parents=True, exist_ok=True)
LOG_PATH.parent.mkdir(parents=True, exist_ok=True)

SUPPORTED_EXT = {".pdf", ".docx", ".pptx"}
Unit = Dict[str, object]

LOGGER = build_script_logger("scan_and_extract", LOG_PATH)


def sha256_file(path: Path, chunk_size: int = 1024 * 1024) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        while True:
            block = handle.read(chunk_size)
            if not block:
                break
            digest.update(block)
    return digest.hexdigest()


def file_stat_info(path: Path) -> Tuple[int, str]:
    stat = path.stat()
    modified = datetime.fromtimestamp(stat.st_mtime, tz=timezone.utc).isoformat()
    return int(stat.st_size), modified


def infer_nature(path: Path) -> str:
    parts = path.parts
    if "V1_RAG" in parts:
        idx = parts.index("V1_RAG")
        if idx + 1 < len(parts):
            return parts[idx + 1]

    try:
        rel = path.relative_to(RAW_DIR)
    except ValueError:
        return ""

    return rel.parts[0] if rel.parts else ""


def is_in_converted(path: Path) -> bool:
    try:
        path.resolve().relative_to(CONVERTED_DIR)
        return True
    except ValueError:
        return False


def map_converted_to_origin(converted_path: Path) -> Optional[Path]:
    try:
        relative = converted_path.resolve().relative_to(CONVERTED_DIR)
    except ValueError:
        return None

    modern_candidate = (RAW_DIR / relative).resolve()
    suffix = converted_path.suffix.lower()

    if suffix == ".docx":
        legacy = modern_candidate.with_suffix(".doc")
        if legacy.exists():
            return legacy
        if modern_candidate.exists():
            return modern_candidate
        return None

    if suffix == ".pptx":
        legacy = modern_candidate.with_suffix(".ppt")
        if legacy.exists():
            return legacy
        if modern_candidate.exists():
            return modern_candidate
        return None

    if modern_candidate.exists():
        return modern_candidate

    return None


def extracted_output_path(source_id: str) -> Path:
    return EXTRACTED_DIR / f"{source_id}.json"


def load_registry() -> List[Row]:
    return load_registry_rows(REGISTRY_PATH)


def write_registry(rows: List[Row]) -> None:
    write_registry_rows(REGISTRY_PATH, rows, logger=LOGGER)


def index_registry(rows: List[Row]) -> Tuple[Dict[str, int], Dict[str, int]]:
    by_path: Dict[str, int] = {}
    by_hash: Dict[str, int] = {}

    for idx, row in enumerate(rows):
        canonical = row.get("canonical_path", "")
        content_hash = row.get("content_hash", "")

        if canonical and canonical not in by_path:
            by_path[canonical] = idx
        if content_hash and content_hash not in by_hash:
            by_hash[content_hash] = idx

    return by_path, by_hash


def scan_corpus() -> List[Path]:
    files: List[Path] = []

    if CONVERTED_DIR.exists():
        for path in CONVERTED_DIR.rglob("*"):
            if path.is_file() and path.suffix.lower() in SUPPORTED_EXT:
                files.append(path.resolve())

    for path in RAW_DIR.rglob("*"):
        if path.is_file() and path.suffix.lower() in SUPPORTED_EXT:
            files.append(path.resolve())

    unique: Dict[str, Path] = {}
    for path in files:
        unique.setdefault(str(path), path)

    return list(unique.values())


def clear_extraction_fields(row: Row) -> None:
    row["extracted_path"] = ""
    row["last_extracted_time"] = ""


def mark_as_new(row: Row) -> None:
    row["extraction_status"] = EXTRACTION_NEW
    row["extraction_error"] = ""
    clear_extraction_fields(row)
    clear_chunk_tracking(row)


def can_restore_extracted(row: Row) -> bool:
    source_id = row.get("source_id", "")
    if not source_id:
        return False

    return extracted_output_path(source_id).exists()


def resolve_origin_path(canonical: Path) -> str:
    if is_in_converted(canonical):
        origin = map_converted_to_origin(canonical)
        return str(origin.resolve()) if origin and origin.exists() else ""

    return str(canonical.resolve())


def update_seen_row(row: Row, canonical: Path, size: int, modified: str, now: str) -> None:
    origin_path = resolve_origin_path(canonical)
    source_for_nature = Path(origin_path) if origin_path else canonical

    row["origin_path"] = origin_path
    row["canonical_path"] = str(canonical.resolve())
    row["filename"] = canonical.name
    row["nature"] = infer_nature(source_for_nature)
    row["ext"] = canonical.suffix.lower()
    row["size_bytes"] = str(size)
    row["modified_time"] = modified
    row["last_seen_time"] = now


def build_new_row(canonical: Path, content_hash: str, size: int, modified: str, now: str) -> Row:
    origin_path = resolve_origin_path(canonical)
    source_for_nature = Path(origin_path) if origin_path else canonical

    return {
        "source_id": content_hash,
        "origin_path": origin_path,
        "canonical_path": str(canonical.resolve()),
        "filename": canonical.name,
        "nature": infer_nature(source_for_nature),
        "ext": canonical.suffix.lower(),
        "size_bytes": str(size),
        "modified_time": modified,
        "content_hash": content_hash,
        "extraction_status": EXTRACTION_NEW,
        "extraction_error": "",
        "extracted_path": "",
        "last_extracted_time": "",
        "last_seen_time": now,
        "chunk_status": "",
        "nb_chunks": "",
        "last_chunk_time": "",
        "chunk_error": "",
        "chunk_config_hash": "",
        "index_status": "",
        "indexed_chunks": "",
        "last_index_time": "",
        "index_error": "",
        "index_config_hash": "",
    }


def ensure_registry_entries(rows: List[Row], corpus_paths: List[Path]) -> List[Row]:
    by_path, by_hash = index_registry(rows)
    now = utc_now_iso()
    seen_canonical: set[str] = set()

    for canonical in corpus_paths:
        canonical = canonical.resolve()
        canonical_str = str(canonical)
        seen_canonical.add(canonical_str)

        size_bytes, modified_time = file_stat_info(canonical)

        existing_by_path = by_path.get(canonical_str)
        if existing_by_path is not None:
            row = rows[existing_by_path]
            previous_hash = row.get("content_hash", "")
            previous_status = row.get("extraction_status", "")
            metadata_changed = (
                row.get("size_bytes", "") != str(size_bytes)
                or row.get("modified_time", "") != modified_time
            )

            update_seen_row(row, canonical, size_bytes, modified_time, now)

            if previous_status == EXTRACTION_MISSING:
                if can_restore_extracted(row):
                    row["extraction_status"] = EXTRACTION_EXTRACTED
                    row["extraction_error"] = ""
                    row["extracted_path"] = str(extracted_output_path(row["source_id"]).resolve())
                    set_chunk_pending(row)
                else:
                    mark_as_new(row)

            if metadata_changed:
                new_hash = sha256_file(canonical)
                if new_hash != previous_hash:
                    row["source_id"] = new_hash
                    row["content_hash"] = new_hash
                    mark_as_new(row)
                elif row.get("extraction_status", "") in ("", EXTRACTION_FAILED):
                    mark_as_new(row)

            if row.get("content_hash", ""):
                by_hash[row["content_hash"]] = existing_by_path
            continue

        content_hash = sha256_file(canonical)
        existing_by_hash = by_hash.get(content_hash)
        if existing_by_hash is not None:
            row = rows[existing_by_hash]
            previous_path = row.get("canonical_path", "")

            update_seen_row(row, canonical, size_bytes, modified_time, now)
            row["source_id"] = content_hash
            row["content_hash"] = content_hash

            if row.get("extraction_status") in ("", EXTRACTION_MISSING):
                if can_restore_extracted(row):
                    row["extraction_status"] = EXTRACTION_EXTRACTED
                    row["extraction_error"] = ""
                    row["extracted_path"] = str(extracted_output_path(row["source_id"]).resolve())
                    set_chunk_pending(row)
                else:
                    mark_as_new(row)
            elif row.get("extraction_status") == EXTRACTION_FAILED:
                mark_as_new(row)

            by_path.pop(previous_path, None)
            by_path[canonical_str] = existing_by_hash
            continue

        rows.append(build_new_row(canonical, content_hash, size_bytes, modified_time, now))
        new_index = len(rows) - 1
        by_path[canonical_str] = new_index
        by_hash[content_hash] = new_index

    for row in rows:
        canonical = row.get("canonical_path", "")
        if canonical and canonical not in seen_canonical:
            row["extraction_status"] = EXTRACTION_MISSING
            row["extraction_error"] = ""
            clear_chunk_tracking(row)

    for row in rows:
        extraction_status = row.get("extraction_status", "")
        if extraction_status == EXTRACTION_EXTRACTED:
            if not row.get("chunk_status", ""):
                set_chunk_pending(row)
            continue
        clear_chunk_tracking(row)

    return rows


def purge_missing_artifacts(rows: List[Row]) -> int:
    removed = 0
    for row in rows:
        if row.get("extraction_status") != EXTRACTION_MISSING:
            continue

        candidates: List[Path] = []
        extracted_path = row.get("extracted_path", "").strip()
        if extracted_path:
            candidates.append(Path(extracted_path))

        source_id = row.get("source_id", "").strip()
        if source_id:
            candidates.append(extracted_output_path(source_id))

        seen: set[str] = set()
        for candidate in candidates:
            candidate_resolved = candidate.resolve()
            candidate_key = str(candidate_resolved)
            if candidate_key in seen:
                continue
            seen.add(candidate_key)

            if not path_is_under(EXTRACTED_DIR, candidate_resolved):
                LOGGER.warning("Skipped purge outside extracted dir: %s", candidate_resolved)
                continue

            if candidate_resolved.exists():
                candidate_resolved.unlink()
                removed += 1

        row["extracted_path"] = ""

    return removed


def normalize_text_block(value: str) -> str:
    text = value.replace("\r\n", "\n").replace("\r", "\n").replace("\x0b", "\n")
    text = text.replace("\u00A0", " ")
    cleaned_chars: List[str] = []
    for char in text:
        codepoint = ord(char)
        if (codepoint < 32 and char not in ("\n", "\r", "\t")) or codepoint == 127:
            cleaned_chars.append(" ")
            continue
        cleaned_chars.append(char)
    text = "".join(cleaned_chars)

    normalized_lines: List[str] = []
    previous_blank = False

    for line in text.split("\n"):
        stripped = line.strip()
        if stripped:
            normalized_lines.append(stripped)
            previous_blank = False
            continue
        if not previous_blank:
            normalized_lines.append("")
            previous_blank = True

    return "\n".join(normalized_lines).strip()


def materialize_units(raw_units: List[Unit]) -> Tuple[str, List[Unit]]:
    text_parts: List[str] = []
    units: List[Unit] = []
    cursor = 0

    for raw_unit in raw_units:
        raw_text = str(raw_unit.get("text", ""))
        text = normalize_text_block(raw_text)
        if not text:
            continue

        if text_parts:
            cursor += 2  # separator inserted between units in the global text

        start_char = cursor
        cursor += len(text)

        unit: Unit = {k: v for k, v in raw_unit.items() if k != "text"}
        unit["text"] = text
        unit["start_char"] = start_char
        unit["end_char"] = cursor
        units.append(unit)
        text_parts.append(text)

    return "\n\n".join(text_parts), units


def extract_pdf_units(path: Path) -> List[Unit]:
    units: List[Unit] = []
    with pymupdf.open(str(path)) as document:
        for page_index, page in enumerate(document, start=1):
            units.append({
                "unit_type": "page",
                "unit_index": page_index,
                "text": page.get_text("text"),
            })
    return units


def extract_docx_units(path: Path) -> List[Unit]:
    document = Document(str(path))
    parsed: List[Tuple[str, bool]] = []
    has_heading = False

    for paragraph in document.paragraphs:
        text = normalize_text_block(paragraph.text or "")
        if not text:
            continue

        style_name = ""
        style_id = ""
        if paragraph.style and paragraph.style.name:
            style_name = paragraph.style.name.strip().lower()
        if paragraph.style and paragraph.style.style_id:
            style_id = paragraph.style.style_id.strip().lower()

        is_heading = style_name.startswith("heading") or style_id.startswith("heading")
        has_heading = has_heading or is_heading
        parsed.append((text, is_heading))

    if not parsed:
        return []

    if not has_heading:
        return [
            {
                "unit_type": "paragraph",
                "unit_index": index,
                "text": text,
            }
            for index, (text, _) in enumerate(parsed, start=1)
        ]

    units: List[Unit] = []
    current_parts: List[str] = []
    current_title = ""
    section_index = 0

    for text, is_heading in parsed:
        if is_heading:
            if current_parts:
                section_index += 1
                section_unit: Unit = {
                    "unit_type": "section",
                    "unit_index": section_index,
                    "text": "\n".join(current_parts),
                }
                if current_title:
                    section_unit["title"] = current_title
                units.append(section_unit)
                current_parts = []

            current_title = text
            current_parts.append(text)
            continue

        current_parts.append(text)

    if current_parts:
        section_index += 1
        section_unit = {
            "unit_type": "section",
            "unit_index": section_index,
            "text": "\n".join(current_parts),
        }
        if current_title:
            section_unit["title"] = current_title
        units.append(section_unit)

    return units


def extract_pptx_units(path: Path) -> List[Unit]:
    presentation = Presentation(str(path))
    units: List[Unit] = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        parts: List[str] = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                value = shape.text.strip()
                if value:
                    parts.append(value)

        if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(notes)

        units.append({
            "unit_type": "slide",
            "unit_index": slide_index,
            "text": "\n".join(parts),
        })

    return units


def extract_units(path: Path) -> List[Unit]:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return extract_pdf_units(path)
    if suffix == ".docx":
        return extract_docx_units(path)
    if suffix == ".pptx":
        return extract_pptx_units(path)
    raise ValueError(f"Unsupported extension for extraction: {suffix}")


def build_extracted_json(row: Row, text: str, units: List[Unit]) -> Dict[str, object]:
    return {
        "text": text,
        "units": units,
        "stats": {
            "unit_count": len(units),
            "char_count": len(text),
        },
        "meta": {
            "source_id": row["source_id"],
            "origin_path": row["origin_path"],
            "canonical_path": row["canonical_path"],
            "filename": row["filename"],
            "nature": row["nature"],
            "ext": row["ext"],
            "modified_time": row["modified_time"],
        },
    }


def run_extraction(rows: List[Row], include_extracted: bool = False) -> List[Row]:
    if include_extracted:
        targets = [
            row
            for row in rows
            if row.get("extraction_status") in (EXTRACTION_NEW, EXTRACTION_EXTRACTED)
        ]
    else:
        targets = [row for row in rows if row.get("extraction_status") == EXTRACTION_NEW]

    LOGGER.info(
        "Extraction started for %d files | include_extracted=%s",
        len(targets),
        include_extracted,
    )

    for row in tqdm(targets, desc="Extraction"):
        canonical = Path(row["canonical_path"])
        try:
            if not canonical.exists():
                row["extraction_status"] = EXTRACTION_MISSING
                row["extraction_error"] = "canonical_path missing on disk"
                clear_chunk_tracking(row)
                continue

            raw_units = extract_units(canonical)
            text, units = materialize_units(raw_units)
            if not text.strip():
                row["extraction_status"] = EXTRACTION_FAILED
                row["extraction_error"] = "empty extracted text"
                clear_chunk_tracking(row)
                continue

            output = extracted_output_path(row["source_id"])
            payload = build_extracted_json(row, text, units)
            output.write_text(
                json.dumps(payload, ensure_ascii=False, indent=2),
                encoding="utf-8",
            )

            row["extracted_path"] = str(output.resolve())
            row["last_extracted_time"] = utc_now_iso()
            row["extraction_status"] = EXTRACTION_EXTRACTED
            row["extraction_error"] = ""
            set_chunk_pending(row)

        except Exception as exc:  # noqa: BLE001
            row["extraction_status"] = EXTRACTION_FAILED
            row["extraction_error"] = str(exc)
            clear_extraction_fields(row)
            clear_chunk_tracking(row)
            LOGGER.exception("Extraction failed for %s", canonical)

    return rows


def extraction_status_counts(rows: List[Row]) -> Dict[str, int]:
    counts: Dict[str, int] = {}
    for row in rows:
        status = row.get("extraction_status", "")
        counts[status] = counts.get(status, 0) + 1

    return counts


def parse_args() -> argparse.Namespace:  # pragma: no cover
    parser = argparse.ArgumentParser(description="Scan files, update registry, extract text.")
    parser.add_argument(
        "--scan-only",
        action="store_true",
        help="Update registry only, without extraction.",
    )
    parser.add_argument(
        "--keep-missing-json",
        action="store_true",
        help="Keep extracted JSON files even if extraction_status becomes MISSING.",
    )
    parser.add_argument(
        "--reextract-extracted",
        action="store_true",
        help="Re-extract files already in EXTRACTED status (useful after schema changes).",
    )
    return parser.parse_args()


def run_scan_and_extract(
    *,
    scan_only: bool,
    keep_missing_json: bool,
    reextract_extracted: bool,
) -> Dict[str, object]:
    LOGGER.info("Run started | RAW_DIR=%s", RAW_DIR)

    rows = load_registry()
    corpus = scan_corpus()
    LOGGER.info("Scan found %d files", len(corpus))

    rows = ensure_registry_entries(rows, corpus)
    purged_before = 0
    if not keep_missing_json:
        purged_before = purge_missing_artifacts(rows)
        if purged_before:
            LOGGER.info("Purged %d extracted JSON files for MISSING sources", purged_before)
    write_registry(rows)

    purged_after = 0
    if not scan_only:
        rows = run_extraction(rows, include_extracted=reextract_extracted)
        if not keep_missing_json:
            purged_after = purge_missing_artifacts(rows)
            if purged_after:
                LOGGER.info("Purged %d extracted JSON files for MISSING sources", purged_after)
        write_registry(rows)

    counts = extraction_status_counts(rows)
    LOGGER.info("Run completed | extraction_status_counts=%s", counts)

    return {
        "extraction_status_counts": counts,
        "purged_before": purged_before,
        "purged_after": purged_after,
        "registry_path": str(REGISTRY_PATH),
        "extracted_dir": str(EXTRACTED_DIR),
        "scan_only": scan_only,
        "keep_missing_json": keep_missing_json,
        "reextract_extracted": reextract_extracted,
    }


def main() -> None:  # pragma: no cover
    args = parse_args()
    result = run_scan_and_extract(
        scan_only=args.scan_only,
        keep_missing_json=args.keep_missing_json,
        reextract_extracted=args.reextract_extracted,
    )

    print("Extraction status:", result["extraction_status_counts"])
    print("Registry:", REGISTRY_PATH)
    print("Extracted:", EXTRACTED_DIR)


if __name__ == "__main__":  # pragma: no cover
    main()
