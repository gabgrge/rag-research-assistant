"""Tests for document extraction drivers, extraction orchestration, and main pipeline run in src/pipeline/scan_and_extract.py."""

import json
from pathlib import Path
from unittest.mock import MagicMock, patch

import pytest
from docx import Document
from pptx import Presentation

from src.pipeline import scan_and_extract as scan_mod


# ============================================================================
# 1. Extraction Drivers Tests
# ============================================================================

class TestExtractionDrivers:
    """Tests for format-specific extractors (extract_pdf_units, extract_docx_units, extract_pptx_units)."""

    def test_extract_pdf_units_reads_pages_via_pymupdf(self, tmp_path):
        """Test extract_pdf_units extracts page units from a PDF using pymupdf."""
        # Arrange
        pdf_file = tmp_path / "sample.pdf"
        pdf_file.touch()

        mock_page_1 = MagicMock()
        mock_page_1.get_text.return_value = "Page 1 Content"
        mock_page_2 = MagicMock()
        mock_page_2.get_text.return_value = "Page 2 Content"

        mock_doc = [mock_page_1, mock_page_2]
        mock_fitz_open = MagicMock()
        mock_fitz_open.return_value.__enter__.return_value = mock_doc

        with patch("pymupdf.open", mock_fitz_open):
            # Act
            units = scan_mod.extract_pdf_units(pdf_file)

        # Assert
        assert len(units) == 2
        assert units[0] == {"unit_type": "page", "unit_index": 1, "text": "Page 1 Content"}
        assert units[1] == {"unit_type": "page", "unit_index": 2, "text": "Page 2 Content"}

    def test_extract_docx_units_reads_paragraphs_without_headings(self, tmp_path):
        """Test extract_docx_units returns paragraph units using a valid saved docx file."""
        # Arrange
        docx_file = tmp_path / "sample.docx"
        doc = Document()
        doc.add_paragraph("Paragraph 1")
        doc.add_paragraph("Paragraph 2")
        doc.save(str(docx_file))

        # Act
        units = scan_mod.extract_docx_units(docx_file)

        # Assert
        assert len(units) == 2
        assert units[0]["unit_type"] == "paragraph"
        assert units[0]["text"] == "Paragraph 1"
        assert units[1]["text"] == "Paragraph 2"

    def test_extract_pptx_units_reads_slides_and_notes(self, tmp_path):
        """Test extract_pptx_units aggregates shape text and speaker notes using a valid saved pptx file."""
        # Arrange
        pptx_file = tmp_path / "sample.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Add text box and speaker notes
        tx_box = slide.shapes.add_textbox(0, 0, 100, 100)
        tx_box.text_frame.text = "Slide 1 Title"
        slide.notes_slide.notes_text_frame.text = "Speaker Notes"

        prs.save(str(pptx_file))

        # Act
        units = scan_mod.extract_pptx_units(pptx_file)

        # Assert
        assert len(units) == 1
        assert units[0]["unit_type"] == "slide"
        assert "Slide 1 Title" in units[0]["text"]
        assert "Speaker Notes" in units[0]["text"]

    def test_extract_units_dispatches_by_extension(self, tmp_path):
        """Test extract_units raises ValueError for unsupported extensions."""
        txt_file = tmp_path / "notes.txt"
        txt_file.touch()

        with pytest.raises(ValueError, match="Unsupported extension"):
            scan_mod.extract_units(txt_file)


# ============================================================================
# 2. Text Normalization & Materialization Tests
# ============================================================================

class TestMaterializeUnits:
    """Tests for text normalization and char offset computation in materialize_units."""

    def test_materialize_units_computes_offsets_and_concatenates_text(self):
        """Test materialize_units sets start_char and end_char while building global text."""
        # Arrange
        raw_units = [
            {"unit_type": "page", "unit_index": 1, "text": "Hello World"},
            {"unit_type": "page", "unit_index": 2, "text": "Second Page"},
        ]

        # Act
        full_text, units = scan_mod.materialize_units(raw_units)

        # Assert
        assert full_text == "Hello World\n\nSecond Page"
        assert len(units) == 2
        assert units[0]["start_char"] == 0
        assert units[0]["end_char"] == 11
        assert units[1]["start_char"] == 13  # 11 + 2 chars for "\n\n"
        assert units[1]["end_char"] == 24


# ============================================================================
# 3. Extraction Run Tests
# ============================================================================

class TestRunExtraction:
    """Tests for run_extraction loop over registry rows."""

    def test_run_extraction_processes_new_rows_and_writes_json(
        self, mock_pipeline_env, sample_registry_row
    ):
        """Test run_extraction extracts NEW rows, writes JSON artifact, and updates row status."""
        # Arrange
        raw_file = mock_pipeline_env["raw_dir"] / "file1.pdf"
        raw_file.touch()

        path_str = str(raw_file.resolve())
        row = sample_registry_row(
            source_id="hash_123",
            origin_path=path_str,
            canonical_path=path_str,
            filename="file1.pdf",
            ext=".pdf",
            extraction_status=scan_mod.EXTRACTION_NEW,
        )

        mock_raw_units = [{"unit_type": "page", "unit_index": 1, "text": "Extracted Content"}]

        with patch.object(scan_mod, "extract_units", return_value=mock_raw_units):
            # Act
            updated_rows = scan_mod.run_extraction([row])

        # Assert
        updated_row = updated_rows[0]
        assert updated_row["extraction_status"] == scan_mod.EXTRACTION_EXTRACTED
        assert updated_row["extraction_error"] == ""

        json_path = Path(updated_row["extracted_path"])
        assert json_path.exists()

        payload = json.loads(json_path.read_text(encoding="utf-8"))
        assert payload["text"] == "Extracted Content"
        assert payload["meta"]["source_id"] == "hash_123"

    def test_run_extraction_handles_failure(
        self, mock_pipeline_env, sample_registry_row
    ):
        """Test run_extraction catches exceptions, marks EXTRACTION_FAILED, and logs error."""
        # Arrange
        raw_file = mock_pipeline_env["raw_dir"] / "bad.pdf"
        raw_file.write_text("Corrupt")

        path_str = str(raw_file.resolve())
        row = sample_registry_row(
            source_id="hash_bad",
            origin_path=path_str,
            canonical_path=path_str,
            filename="bad.pdf",
            ext=".pdf",
            extraction_status=scan_mod.EXTRACTION_NEW,
        )

        with patch.object(scan_mod, "extract_units", side_effect=RuntimeError("Parsing error")):
            # Act
            updated_rows = scan_mod.run_extraction([row])

        # Assert
        updated_row = updated_rows[0]
        assert updated_row["extraction_status"] == scan_mod.EXTRACTION_FAILED
        assert "Parsing error" in updated_row["extraction_error"]


# ============================================================================
# 4. Orchestration Tests
# ============================================================================

class TestRunScanAndExtract:
    """Tests for run_scan_and_extract orchestrator function."""

    def test_run_scan_and_extract_end_to_end(self, mock_pipeline_env):
        """Test run_scan_and_extract full execution path."""
        # Arrange
        doc = mock_pipeline_env["raw_dir"] / "doc.pdf"
        doc.write_text("Hello PDF")

        mock_raw_units = [{"unit_type": "page", "unit_index": 1, "text": "Hello PDF"}]

        with patch.object(scan_mod, "extract_units", return_value=mock_raw_units):
            # Act
            result = scan_mod.run_scan_and_extract(
                scan_only=False,
                keep_missing_json=False,
                reextract_extracted=False,
            )

        # Assert
        assert result["scan_only"] is False
        assert result["extraction_status_counts"].get(scan_mod.EXTRACTION_EXTRACTED) == 1

        # Verify registry content on disk
        rows = scan_mod.load_registry()
        assert len(rows) == 1
        assert rows[0]["extraction_status"] == scan_mod.EXTRACTION_EXTRACTED
